from fastapi import FastAPI
import uvicorn
import re
import moviepy.editor as mp
from moviepy.video.tools.subtitles import SubtitlesClip
import srt
import requests
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches
import datetime
import subprocess
from fastapi.middleware.cors import CORSMiddleware
import spire.presentation
from fastapi.staticfiles import StaticFiles
#from ppt_to_png import *
app = FastAPI()
app.mount("/output", StaticFiles(directory="output"), name="output")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

#Ai润色
def text_polish(texts):
    api="sk-97911fd3236c466f9b4d65dfddef6385"
    client = OpenAI(api_key=api, base_url="https://api.deepseek.com")
    for i,text in enumerate(texts):
        response = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": "You are a helpful assistant"},
                {"role": "user", "content": "给我润色一下这段话，不要产生任何与原文本无关的内容，且最多比原文本多20个字：{}".format(text)},
            ],
            temperature=1.3,
            stream=False
        )
        result=response.choices[0].message.content
        print("润色后文本{}：{}".format(i,response.choices[0].message.content))
        texts[i]=result

    return texts

#读取每一页的备注，返回一个列表
def read_ppt_notes(ppt_path):
    ppt = Presentation(ppt_path)
    notes_texts = []

    for slide in ppt.slides:
        # 默认备注为空
        notes_text = ""
        # 读取备注页的文本
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text
            notes_text = re.sub(r'[^\u4e00-\u9fa5，。！？、；：“”‘’（）【】《》…]', '', notes_text)  # 只保留汉字和标点

        notes_texts.append(notes_text)

    return notes_texts
#删除ppt里不是中文的部分
def remove_non_chinese(text):
    result = re.sub(r'[^\u4e00-\u9fa5，。！？、；：“”‘’（）【】《》…]', '', text)
    return result

# 1. 提取PPT中的文字，返回一个列表，每个元素就是一页ppt里的文字
def extract_text_from_ppt(ppt_path):
    """
    从PPT文件中逐页提取文字内容。
    """
    prs = Presentation(ppt_path)
    texts = []
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text
        texts.append(slide_text.strip())
    for i in range(0,len(texts)):#删除汉字部分
        texts[i]=remove_non_chinese(texts[i])
    print("ppt文本：",texts)
    return texts

# 2. 给每一页ppt添加备注。接收一个列表，列表的每个元素就是一页ppt的备注
def add_note(path,texts):
    # 打开PPT文件
    ppt = Presentation(path)
    # 遍历所有幻灯片，给每张幻灯片添加备注
    for i, slide in enumerate(ppt.slides):
        # 获取或创建备注页
        notes_slide = slide.notes_slide

        # 设置备注内容
        notes_slide.notes_text_frame.text = texts[i]

    # 保存修改后的 PPT
    ppt.save(path)
    print("备注已添加到PPT中！{}".format(path))
# 2. 调用So-VITS API生成语音，每页ppt的文字单独一个音频，然后返回一个包含这些音频文件名的列表
def text_to_speech(texts,reference_audio="man_1.wav",top_k=5,top_p=1,temperature=1):#接收一个字符串列表
    name_list=[];
    url = "http://127.0.0.1:9880/tts"
    for text in texts:
        params = {
            "text": text,
            "text_lang": "zh",
            "ref_audio_path": reference_audio,
            "prompt_lang": "zh",
            "media_type": "wav",
            "top_k":top_k,
            "top_p":top_p,
            "temperature":temperature
        }

        response = requests.get(url, params=params)
        if response.status_code == 200:
            with open("output{}.wav".format(texts.index(text)), "wb") as f:
                name_list.append("output{}.wav".format(texts.index(text)))
                f.write(response.content)
            print("语音合成成功！")
        else:
            print("错误:", response.json())
    print("音频列表：",name_list)
    return name_list
def add_audio_to_ppt(ppt_path, output_path,audio_paths):
    # 加载或创建演示文稿对象
    prs = Presentation(ppt_path)  # 如果文件不存在，请确保创建新演示文稿
    for i,audio_path in enumerate(audio_paths):
        # 获取第一页幻灯片（如果存在）
        slide = prs.slides[i]
        # 添加音频文件
        left = Inches(1)
        top = Inches(1)
        width = Inches(2)
        height = Inches(2)

        # 将音频嵌入第一页幻灯片
        slide.shapes.add_movie(audio_path, left, top, width, height, mime_type='audio/wav')

    # 保存演示文稿
    prs.save(output_path)

def ppt_to_png(ppt_path,output_folder):
    img_file=[]
    # 创建一个Presentation对象
    presentation = spire.presentation.Presentation()
    # 从文件加载名为"输入文档.pptx"的演示文稿数据
    presentation.LoadFromFile(ppt_path)
    # 遍历演示文稿中的每个幻灯片
    for i, slide in enumerate(presentation.Slides):
        # 构建输出文件名，格式为"Output/ToImage_序号.png"
        fileName = output_folder + str(i) + ".png"
        img_file.append(fileName)
        # 将当前幻灯片保存为图像
        image = slide.SaveAsImage()
        # 将图像保存到指定文件名
        image.Save(fileName)
        # 释放图像资源
        #image.Dispose()
    presentation.Dispose()
    return img_file


# 5. FastAPI 端点
@app.post("/process_ppt")
def process_ppt(ppt_path: str, output_path: str,reference_audio: str,choice:int,top_k:int=5,top_p:int=1,temperature:int=1):

    #第一种方案
    if choice==1:
        # 从ppt中提取文本
        texts = extract_text_from_ppt(ppt_path)
        # 添加备注
        add_note(ppt_path, texts)
        # 文本转语音
        audio_files = text_to_speech(texts, reference_audio, top_k, top_p, temperature)
        # 往ppt里插入音频
        add_audio_to_ppt(ppt_path, output_path, audio_files)

    if choice == 2:
        # 从ppt中提取文本
        texts = extract_text_from_ppt(ppt_path)
        #AI润色文本
        texts = text_polish(texts)
        # 添加备注
        add_note(ppt_path, texts)
        # 文本转语音
        audio_files = text_to_speech(texts, reference_audio, top_k, top_p, temperature)
        # 往ppt里插入音频
        add_audio_to_ppt(ppt_path, output_path, audio_files)

    if choice == 3:
        #读取备注
        texts=read_ppt_notes(ppt_path)
        # 文本转语音
        audio_files = text_to_speech(texts, reference_audio, top_k, top_p, temperature)
        # 往ppt里插入音频
        add_audio_to_ppt(ppt_path, output_path, audio_files)
    return {"message": "PPT 处理完成", "output_path": output_path}

# 运行 FastAPI
if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
