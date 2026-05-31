from fastapi import FastAPI, Body
import uvicorn
import re
import moviepy.editor as mp
from moviepy.video.tools.subtitles import SubtitlesClip
import srt
import requests
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches
import datetime
import subprocess
from fastapi.middleware.cors import CORSMiddleware
import spire.presentation
import os
from pathlib import Path
import logging
from fastapi.staticfiles import StaticFiles
# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ppt_processor.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("ppt_processor")

# 创建FastAPI应用
app = FastAPI()
app.mount("/output", StaticFiles(directory="output"), name="output")
# 添加CORS中间件，允许跨域请求
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)


# Ai润色函数 - 使用DeepSeek API提升文本质量
def text_polish(texts):
    """
    使用DeepSeek API对文本进行润色

    参数:
    texts (list): 需要润色的文本列表

    返回:
    list: 润色后的文本列表
    """
    try:
        # 注意: 建议将API密钥存储在环境变量中，而不是硬编码在代码中
        api = "sk-97911fd3236c466f9b4d65dfddef6385"
        client = OpenAI(api_key=api, base_url="https://api.deepseek.com")

        for i, text in enumerate(texts):
            if not text or text.strip() == "":
                logger.warning(f"文本{i}为空，跳过润色")
                continue

            try:
                response = client.chat.completions.create(
                    model="deepseek-chat",
                    messages=[
                        {"role": "system", "content": "You are a helpful assistant"},
                        {"role": "user",
                         "content": f"给我润色一下这段话，不要产生任何与原文本无关的内容，且最多比原文本多20个字：{text}"},
                    ],
                    temperature=1.3,
                    stream=False
                )
                result = response.choices[0].message.content
                logger.info(f"润色后文本{i}：{result}")
                texts[i] = result
            except Exception as e:
                logger.error(f"文本{i}润色失败: {e}")
                # 失败时保留原文本

        return texts
    except Exception as e:
        logger.error(f"文本润色过程发生错误: {e}")
        return texts  # 发生错误时返回原始文本


# 读取每一页的备注，返回一个列表
def read_ppt_notes(ppt_path):
    """
    读取PPT中每一页的备注

    参数:
    ppt_path (str): PPT文件路径

    返回:
    list: 包含每页备注文本的列表
    """
    try:
        logger.info(f"开始读取PPT备注: {ppt_path}")
        ppt = Presentation(ppt_path)
        notes_texts = []

        for i, slide in enumerate(ppt.slides):
            # 默认备注为空
            notes_text = ""
            # 读取备注页的文本
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text
                notes_text = re.sub(r'[^\u4e00-\u9fa5，。！？、；：""''（）【】《》…]', '', notes_text)  # 只保留汉字和标点

            notes_texts.append(notes_text)
            logger.debug(f"第{i + 1}页备注: {notes_text}")

        logger.info(f"成功读取{len(notes_texts)}页PPT备注")
        return notes_texts
    except Exception as e:
        logger.error(f"读取PPT备注失败: {e}")
        raise


# 删除文本中非中文的部分
def remove_non_chinese(text):
    """
    删除文本中非中文字符

    参数:
    text (str): 输入文本

    返回:
    str: 只保留中文和指定标点的文本
    """
    result = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9，。！？、；：""''（）《》【】]', '', text)
    return result


# 提取PPT中的文字，返回一个列表，每个元素就是一页ppt里的文字
def extract_text_from_ppt(ppt_path):
    """
    从PPT文件中逐页提取文字内容

    参数:
    ppt_path (str): PPT文件路径

    返回:
    list: 包含每页文本内容的列表
    """
    try:
        logger.info(f"开始从PPT提取文本: {ppt_path}")
        prs = Presentation(ppt_path)
        texts = []

        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text
            texts.append(slide_text.strip())

        # 处理文本，只保留中文部分
        for i in range(len(texts)):
            texts[i] = remove_non_chinese(texts[i])

        logger.info(f"成功从PPT提取{len(texts)}页文本")
        logger.debug(f"提取的文本: {texts}")
        return texts
    except Exception as e:
        logger.error(f"从PPT提取文本失败: {e}")
        raise


# 给每一页ppt添加备注
def add_note(path, texts):
    """
    给PPT每页添加备注

    参数:
    path (str): PPT文件路径
    texts (list): 备注文本列表
    """
    try:
        logger.info(f"开始添加备注到PPT: {path}")
        # 打开PPT文件
        ppt = Presentation(path)

        # 检查文本列表长度与幻灯片数量是否匹配
        if len(texts) > len(ppt.slides):
            logger.warning(f"备注数量({len(texts)})大于幻灯片数量({len(ppt.slides)})")

        # 遍历所有幻灯片，给每张幻灯片添加备注
        for i, slide in enumerate(ppt.slides):
            if i >= len(texts):
                logger.warning(f"没有为第{i + 1}页幻灯片提供备注")
                break

            # 获取或创建备注页
            notes_slide = slide.notes_slide
            # 设置备注内容
            notes_slide.notes_text_frame.text = texts[i]
            logger.debug(f"已添加备注到第{i + 1}页: {texts[i][:30]}...")

        # 保存修改后的 PPT
        ppt.save(path)
        logger.info(f"备注已成功添加到PPT中: {path}")
    except Exception as e:
        logger.error(f"添加备注到PPT失败: {e}")
        raise


# 文本转语音函数 - 单个句子处理
def text_to_speech_sentence(text, index, sub_index, reference_audio, top_k=5, top_p=1, temperature=1):
    """
    将单个文本转换为语音

    参数:
    text (str): 需要转换的文本
    index (int): 幻灯片索引
    sub_index (int): 句子在幻灯片中的索引
    reference_audio (str): 参考音频文件
    top_k (int): TTS参数
    top_p (float): TTS参数
    temperature (float): TTS参数

    返回:
    str: 生成的音频文件路径，失败时返回None
    """
    # 打印全部参数
    logger.info(f"转换文本参数: text={text}, index={index}, sub_index={sub_index}, reference_audio={reference_audio}, top_k={top_k}, top_p={top_p}, temperature={temperature}")
    # 检查文本是否为空
    if not text or text.strip() == "":
        logger.warning(f"第{index + 1}页第{sub_index + 1}句文本为空，跳过生成")
        return None

    url = "http://127.0.0.1:9880/tts"
    output_file = f"output/output_{index}_{sub_index}.wav"

    # 设置TTS API参数
    params = {
        "text": text,
        "text_lang": "zh",
        "ref_audio_path": reference_audio,
        "prompt_lang": "zh",
        "media_type": "wav",
        "top_k": top_k,
        "top_p": top_p,
        "temperature": temperature
    }

    try:
        # 发送请求到TTS服务
        logger.info(f"请求TTS API转换文本: {text[:30]}...")
        logger.info(f"请求TTS API转换文本 url: {url} {params}")
        response = requests.get(url, params=params)

        # 检查响应状态
        if response.status_code == 200:
            # 保存音频文件
            with open(output_file, "wb") as f:
                f.write(response.content)
            logger.info(f"已生成句子音频: {output_file}")

            return output_file
        else:
            logger.error(f"错误响应 ({response.status_code}): {response.text}")
            return None
    except requests.exceptions.RequestException as e:
        logger.error(f"TTS API请求错误: {e}")
        return None
    except Exception as e:
        logger.error(f"生成音频过程中发生错误: {e}")
        return None


# 为PPT的每一页生成语音
def text_to_speech_for_slides(texts, reference_audio, top_k=5, top_p=1, temperature=1):
    """
    为PPT的每一页文本生成语音

    参数:
    texts (list): 文本列表，每个元素对应一页PPT
    reference_audio (str): 参考音频
    top_k (int): TTS参数
    top_p (float): TTS参数
    temperature (float): TTS参数

    返回:
    list: 生成的音频文件路径列表
    """
    logger.info("开始为PPT生成语音")
    audio_files = []

    for i, text in enumerate(texts):
        if not text or text.strip() == "":
            logger.warning(f"第{i + 1}页文本为空，跳过生成")
            audio_files.append(None)
            continue

        # 为每页生成音频
        audio_file = text_to_speech_sentence(text, i, 0, reference_audio, top_k, top_p, temperature)
        audio_files.append(audio_file)

    logger.info(f"已完成{len(audio_files)}页PPT的语音生成")
    return audio_files


# 往PPT添加音频
def add_audio_to_ppt(ppt_path, output_path, audio_paths):
    """
    将音频添加到PPT中

    参数:
    ppt_path (str): 输入PPT路径
    output_path (str): 输出PPT路径
    audio_paths (list): 音频文件路径列表
    """
    try:
        logger.info(f"开始将音频添加到PPT: {ppt_path}")
        # 确保目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        prs = Presentation(ppt_path)

        # 检查音频文件数量
        if len(audio_paths) > len(prs.slides):
            logger.warning(f"音频文件数量({len(audio_paths)})大于幻灯片数量({len(prs.slides)})")

        audio_added = 0
        for i, audio_path in enumerate(audio_paths):
            # 检查索引是否有效
            if i >= len(prs.slides):
                logger.warning(f"音频索引{i}超出了幻灯片数量范围，跳过")
                continue

            # 检查音频文件是否存在
            if not audio_path or not os.path.exists(audio_path):
                logger.warning(f"音频文件不存在: {audio_path}，跳过第{i + 1}页")
                continue

            # 获取幻灯片
            slide = prs.slides[i]

            # 设置音频位置
            left = Inches(1)
            top = Inches(1)
            width = Inches(2)
            height = Inches(2)

            # 添加音频到幻灯片
            slide.shapes.add_movie(audio_path, left, top, width, height, mime_type='audio/wav')
            audio_added += 1

        # 保存PPT
        prs.save(output_path)
        logger.info(f"已成功添加{audio_added}个音频到PPT: {output_path}")
    except Exception as e:
        logger.error(f"添加音频到PPT时发生错误: {e}")
        raise


# 将PPT转换为PNG图片
def ppt_to_png(ppt_path, output_folder):
    """
    将PPT转换为PNG图片

    参数:
    ppt_path (str): PPT文件路径
    output_folder (str): 输出文件夹

    返回:
    list: 生成的图片文件路径列表
    """
    try:
        logger.info(f"开始将PPT转换为PNG: {ppt_path}")

        # 确保输出目录存在
        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        img_file = []
        # 创建Presentation对象
        presentation = spire.presentation.Presentation()

        # 加载PPT文件
        presentation.LoadFromFile(ppt_path)

        # 遍历PPT中的每张幻灯片
        for i, slide in enumerate(presentation.Slides):
            # 构建输出文件名
            fileName = output_folder + str(i) + ".png"
            img_file.append(fileName)

            # 保存为图片
            image = slide.SaveAsImage()
            image.Save(fileName)
            logger.debug(f"已保存第{i + 1}页为图片: {fileName}")

        # 释放资源
        presentation.Dispose()

        logger.info(f"已成功将PPT转换为{len(img_file)}张PNG图片")
        return img_file
    except Exception as e:
        logger.error(f"PPT转换为PNG失败: {e}")
        raise


# 使用spire.presentation库添加音频到PPT的改进版本
def add_audio_to_ppt_spire_fixed(ppt_path, output_path, audio_paths):
    """
    使用spire.presentation库添加音频到PPT（修复版）

    参数:
    ppt_path (str): 输入PPT路径
    output_path (str): 输出PPT路径
    audio_paths (list): 音频文件路径列表
    """
    try:
        logger.info(f"开始使用Spire库将音频添加到PPT: {ppt_path}")

        # 确保所有路径都是绝对路径
        abs_audio_paths = [os.path.abspath(path) for path in audio_paths if path and os.path.exists(path)]
        abs_ppt_path = os.path.abspath(ppt_path)
        abs_output_path = os.path.abspath(output_path)

        # 创建演示文稿对象
        presentation = spire.presentation.Presentation()

        # 加载PPT文件
        presentation.LoadFromFile(abs_ppt_path)

        # 检查库版本并调整方法调用
        audio_added = 0

        for i, audio_path in enumerate(abs_audio_paths):
            # 确保索引有效
            if i >= presentation.Slides.Count:
                logger.warning(f"音频索引{i}超出了幻灯片数量范围({presentation.Slides.Count})，跳过")
                continue

            # 获取幻灯片
            slide = presentation.Slides[i]

            try:
                # 尝试多种可能的方法添加音频
                try:
                    # 方法1：尝试使用AppendAudioMedia（某些版本中的名称）
                    audio = slide.Shapes.AppendAudioMedia(audio_path, 50, 50, 100, 30)
                    logger.info(f"使用AppendAudioMedia方法添加音频成功")
                except AttributeError:
                    try:
                        # 方法2：尝试使用AddAudioMedia
                        audio = slide.Shapes.AddAudioMedia(audio_path, 50, 50, 100, 30)
                        logger.info(f"使用AddAudioMedia方法添加音频成功")
                    except AttributeError:
                        try:
                            # 方法3：尝试使用AddMediaObject
                            audio = slide.Shapes.AddMediaObject(audio_path, 50, 50, 100, 30)
                            logger.info(f"使用AddMediaObject方法添加音频成功")
                        except AttributeError:
                            # 方法4：查看是否有其他可能的方法
                            methods = [method for method in dir(slide.Shapes) if
                                       'media' in method.lower() or 'audio' in method.lower()]
                            logger.info(f"可用的媒体相关方法: {methods}")

                            # 如果找不到合适的方法，记录错误并继续
                            logger.error(f"找不到适合添加音频的方法，可用方法: {methods}")
                            continue

                # 尝试设置音频自动播放（如果API支持）
                try:
                    if hasattr(audio, 'PlayMode'):
                        audio.PlayMode = spire.presentation.AudioPlayMode.Auto
                    elif hasattr(audio, 'PlaySettings'):
                        audio.PlaySettings.PlayOnClick = True
                        audio.PlaySettings.PlayOnEntry = True
                    logger.info(f"成功设置音频自动播放")
                except Exception as e:
                    logger.warning(f"设置音频自动播放失败: {e}")

                audio_added += 1
                logger.info(f"已添加音频到第{i + 1}页: {audio_path}")
            except Exception as e:
                logger.error(f"向第{i + 1}页添加音频时出错: {e}")

        # 保存PPT到新文件
        # 尝试不同的文件格式常量
        try:
            presentation.SaveToFile(abs_output_path, spire.presentation.FileFormat.Pptx)
        except AttributeError:
            try:
                # 另一种可能的格式枚举
                presentation.SaveToFile(abs_output_path, spire.presentation.FileFormat.PPTX)
            except AttributeError:
                # 如果没有格式常量，直接保存
                presentation.SaveToFile(abs_output_path)

        presentation.Dispose()
        logger.info(f"已成功添加{audio_added}个音频到PPT: {output_path}")

        # 如果没有成功添加任何音频，抛出异常，以便后续可以尝试备用方法
        if audio_added == 0:
            raise Exception("Spire库未能添加任何音频")

        return True
    except Exception as e:
        logger.error(f"使用Spire添加音频到PPT时发生错误: {e}")
        raise


# 创建一个包含音频链接的新PPT
def create_ppt_with_audio_links(ppt_path, output_path, audio_paths):
    """
    创建一个包含音频链接的新PPT

    参数:
    ppt_path (str): 源PPT路径（用于读取内容）
    output_path (str): 输出PPT路径
    audio_paths (list): 音频文件路径列表
    """
    try:
        logger.info(f"开始创建带音频链接的PPT: {output_path}")

        # 读取源PPT
        source_ppt = Presentation(ppt_path)

        # 创建新的PPT
        new_ppt = Presentation()

        # 复制幻灯片并添加音频链接
        audio_added = 0
        valid_audio_paths = [path for path in audio_paths if path and os.path.exists(path)]

        for i, slide in enumerate(source_ppt.slides):
            # 复制幻灯片布局
            blank_slide_layout = new_ppt.slide_layouts[6] if len(new_ppt.slide_layouts) > 6 else new_ppt.slide_layouts[
                0]  # 使用空白布局或默认布局
            new_slide = new_ppt.slides.add_slide(blank_slide_layout)

            # 复制形状
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    txBox = new_slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.text = shape.text

            # 添加音频链接（如果有）
            if i < len(valid_audio_paths):
                audio_path = os.path.abspath(valid_audio_paths[i])

                # 在幻灯片底部添加文本框
                audio_box = new_slide.shapes.add_textbox(
                    Inches(1), Inches(5), Inches(6), Inches(0.5))
                audio_frame = audio_box.text_frame
                p = audio_frame.paragraphs[0]

                # 添加链接
                run = p.add_run()
                run.text = f"点击播放音频 {i + 1}"

                # 添加超链接
                run.hyperlink.address = audio_path

                audio_added += 1
                logger.debug(f"已添加音频链接到第{i + 1}页: {audio_path}")

        # 保存新PPT
        new_ppt.save(output_path)
        logger.info(f"已成功创建带有{audio_added}个音频链接的PPT: {output_path}")

    except Exception as e:
        logger.error(f"创建带音频链接的PPT时发生错误: {e}")
        raise


# PPT转视频端点
@app.post("/ppt_to_video")
def ppt_to_video(ppt_path: str = Body(...), video_output: str = Body(...), video_src_output: str = Body(...), 
                reference_audio: str = Body(...), top_k: int = Body(5), top_p: float = Body(1.0), 
                temperature: float = Body(1.0)):
    """
    将PPT转换为视频

    参数:
    ppt_path (str): PPT文件路径
    video_output (str): 输出视频路径
    video_src_output (str): 带字幕的输出视频路径
    reference_audio (str): 参考音频文件路径，默认为"reference.wav"
    top_k (int): TTS参数，控制采样的多样性，默认为5
    top_p (float): TTS参数，控制采样的累积概率阈值，默认为1.0
    temperature (float): TTS参数，控制采样的随机性，默认为1.0

    返回:
    dict: 处理结果
    """
    try:
        logger.info(f"开始PPT转视频处理: {ppt_path}")
        logger.info(f"TTS参数: reference_audio={reference_audio}, top_k={top_k}, top_p={top_p}, temperature={temperature}")

        # === 1. 检查文件是否存在 ===
        if not os.path.exists(ppt_path):
            return {"error": f"PPT文件不存在: {ppt_path}"}

        # 检查参考音频文件是否存在
        if not os.path.exists(reference_audio):
            logger.warning(f"参考音频文件不存在: {reference_audio}，将使用默认语音")

        # === 2. PPT 转图片 ===
        image_files = ppt_to_png(ppt_path, "ppt_png/")
        logger.info("PPT转PNG成功")

        # === 3. 提取文本并处理 ===
        texts = extract_text_from_ppt(ppt_path)
        logger.info("文本提取成功")

        # 分割文本为句子的辅助函数
        def split_into_sentences(text):
            # 定义句子结束的标点
            sentence_ends = ['。', '！', '？', '；', '.', '!', '?', ';']
            sentences = []
            current_sentence = ""

            for char in text:
                current_sentence += char
                if char in sentence_ends:
                    if current_sentence.strip():
                        sentences.append(current_sentence.strip())
                    current_sentence = ""

            # 处理最后可能没有结束标点的句子
            if current_sentence.strip():
                sentences.append(current_sentence.strip())

            # 如果没有分出句子但有文本，将整个文本作为一句
            if not sentences and text.strip():
                sentences = [text.strip()]

            return sentences

        # 清理文本函数，保留中文、数字、英文字母和标点
        def clean_text(text):
            return re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9，。！？、；：""''（）《》【】,.!?;:()[\]]+', '', text)

        # 处理每个PPT页面的文本，生成音频和字幕
        all_slide_data = []

        for i, text in enumerate(texts):
            cleaned_text = clean_text(text)
            sentences = split_into_sentences(cleaned_text)

            slide_data = {
                "image": image_files[i] if i < len(image_files) else None,
                "sentences": [],
                "total_duration": 0
            }

            # 检查图片是否存在
            if not slide_data["image"] or not os.path.exists(slide_data["image"]):
                logger.warning(f"第{i + 1}页图片不存在，跳过")
                continue

            # 为每个句子生成音频
            for j, sentence in enumerate(sentences):
                if not sentence or sentence.strip() == "":
                    logger.warning(f"第{i + 1}页第{j + 1}句为空，跳过")
                    continue

                audio_file = text_to_speech_sentence(sentence, i, j, reference_audio, top_k, top_p, temperature)

                if audio_file and os.path.exists(audio_file):
                    # 获取音频时长
                    try:
                        audio_clip = mp.AudioFileClip(audio_file)
                        duration = audio_clip.duration
                        audio_clip.close()

                        slide_data["sentences"].append({
                            "text": sentence,
                            "audio": audio_file,
                            "duration": duration
                        })
                        slide_data["total_duration"] += duration
                    except Exception as e:
                        logger.error(f"处理音频文件失败: {audio_file}, 错误: {e}")

            # 只有当有句子成功生成音频时，才添加幻灯片数据
            if slide_data["sentences"]:
                all_slide_data.append(slide_data)

        logger.info("所有句子音频生成成功")

        # === 4. 生成字幕文件 ===
        # 改进的文本分行处理函数
        def split_line(text):
            """更精确的文本分行函数"""
            # 如果文本较短，直接返回
            if len(text) <= 20:
                return [text]

            # 常用标点符号
            punctuation = ['。', '，', '；', '、', '：', '！', '？', '.', ',', ';', ':', '!', '?']

            # 强制分行的最大长度
            max_line_length = 25

            # 第一次尝试：在标点处分行
            lines = []
            current_line = ""

            for char in text:
                current_line += char
                # 当遇到标点且当前行长度至少15个字符，就在此处分行
                if char in punctuation and len(current_line) >= 15:
                    lines.append(current_line)
                    current_line = ""

            # 添加最后一部分
            if current_line:
                lines.append(current_line)

            # 第二次处理：如果第一次分行后有行超过最大长度，再次分行
            final_lines = []
            for line in lines:
                if len(line) > max_line_length:
                    # 尝试在次优先标点处分行
                    sublines = []
                    subline = ""
                    found_break = False

                    for char in line:
                        subline += char
                        if len(subline) >= max_line_length // 2 and char in punctuation:
                            sublines.append(subline)
                            subline = ""
                            found_break = True
                            break

                    if subline:
                        sublines.append(subline)

                    # 如果没找到合适的分割点，按固定长度分割
                    if not found_break:
                        for i in range(0, len(line), max_line_length):
                            sublines.append(line[i:i + max_line_length])

                    final_lines.extend(sublines)
                else:
                    final_lines.append(line)

            # 如果没有分行成功，则强制按长度分行
            if not final_lines:
                for i in range(0, len(text), max_line_length):
                    final_lines.append(text[i:i + max_line_length])

            # 最多返回两行
            return final_lines[:2]

        # 生成带有精确时间的字幕
        def generate_subtitles_with_precise_timing(slide_data, subtitles_file):
            subtitles = []
            subtitle_index = 1
            current_time = datetime.timedelta(seconds=0)

            for slide in slide_data:
                for sentence_data in slide["sentences"]:
                    sentence = sentence_data["text"]

                    # 改进的分行处理
                    lines = split_line(sentence)
                    content = "\n".join(lines)

                    # 计算结束时间
                    end_time = current_time + datetime.timedelta(seconds=sentence_data["duration"])

                    subtitle = srt.Subtitle(
                        index=subtitle_index,
                        start=current_time,
                        end=end_time,
                        content=content
                    )
                    subtitles.append(subtitle)
                    subtitle_index += 1
                    current_time = end_time

            # 写入SRT文件
            with open(subtitles_file, "w", encoding="utf-8") as f:
                f.write(srt.compose(subtitles))

        # 生成字幕文件
        subtitles_file = "ppt_srt.srt"
        generate_subtitles_with_precise_timing(all_slide_data, subtitles_file)
        logger.info("字幕生成成功")

        # === 5. 生成视频 ===
        def create_video_with_precise_timing(slide_data, output_video):
            if not slide_data:
                raise ValueError("没有有效的幻灯片数据用于创建视频")

            clips = []

            for slide in slide_data:
                # 检查图片是否存在
                if not os.path.exists(slide["image"]):
                    logger.warning(f"图片不存在: {slide['image']}，跳过")
                    continue

                if slide["total_duration"] <= 0:
                    logger.warning(f"幻灯片持续时间为0: {slide['image']}，跳过")
                    continue

                # 创建图片片段，持续时间为所有句子音频的总时长
                img_clip = mp.ImageClip(slide["image"], duration=slide["total_duration"])

                # 合并该页的所有句子音频
                audio_clips = []
                for sentence in slide["sentences"]:
                    if os.path.exists(sentence["audio"]):
                        audio_clips.append(mp.AudioFileClip(sentence["audio"]))

                if audio_clips:
                    combined_audio = mp.concatenate_audioclips(audio_clips)
                    img_clip = img_clip.set_audio(combined_audio)

                clips.append(img_clip)

            if not clips:
                raise ValueError("没有有效的视频片段用于合并")

            # 合并所有视频片段
            final_video = mp.concatenate_videoclips(clips, method="compose")
            # 保存视频
            output_dir = os.path.dirname(output_video)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)

            final_video.write_videofile(output_video, codec="libx264", fps=1)

        # 创建视频
        create_video_with_precise_timing(all_slide_data, video_output)
        logger.info("视频合成成功")

        # 添加字幕
        def add_subtitles(video_path, srt_path, output_path):
            # 确保输出目录存在
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)

            command = [
                "ffmpeg", "-i", video_path, "-vf",
                f"subtitles={srt_path}:force_style='FontName=SimHei,FontSize=10,PrimaryColour=&HFFFFFF&,Alignment=2,BorderStyle=3,Outline=1,Shadow=0,BackColour=&H00000000'",
                "-c:a", "copy",
                "-max_muxing_queue_size", "1024",
                output_path, "-y"
            ]

            try:
                result = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                if result.returncode != 0:
                    logger.error(f"添加字幕失败: {result.stderr.decode('utf-8', errors='ignore')}")
                    raise Exception("ffmpeg命令执行失败")
            except Exception as e:
                logger.error(f"执行ffmpeg命令时出错: {e}")
                raise

        # 添加字幕到视频
        add_subtitles(video_output, subtitles_file, video_src_output)
        logger.info("字幕添加成功")

        return {"message": "PPT to video 处理完成", "output_path": video_src_output}
    except Exception as e:
        logger.error(f"PPT转视频处理失败: {e}")
        return {"error": str(e)}

# 运行 FastAPI
if __name__ == "__main__":
    logger.info("启动PPT处理服务...")
    uvicorn.run(app, host="127.0.0.1", port=8000)